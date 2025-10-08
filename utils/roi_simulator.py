from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, Any, List

import pandas as pd

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


@dataclass
class BaselineMetrics:
    spend: float
    revenue: float
    roas: float
    conversions: float
    cpa: float
    impressions: float
    clicks: float
    ctr: float


@dataclass
class ROIReportRequest:
    brand_name: str
    start_date: datetime
    end_date: datetime
    scenario_name: str = "ROI 模擬情境"
    analyst: Optional[str] = None
    notes: Optional[List[str]] = None


def _filter_by_date(df: pd.DataFrame, start: datetime, end: datetime) -> pd.DataFrame:
    if "開始" not in df.columns:
        return df
    mask = (df["開始"] >= pd.Timestamp(start)) & (df["開始"] <= pd.Timestamp(end))
    filtered = df[mask]
    if filtered.empty:
        return df
    return filtered


def calculate_baseline(df: pd.DataFrame, start: datetime, end: datetime) -> BaselineMetrics:
    data = _filter_by_date(df, start, end).copy()

    spend = float(data.get("花費金額 (TWD)", pd.Series(dtype=float)).sum())
    roas_series = data.get("購買 ROAS（廣告投資報酬率）", pd.Series(dtype=float))
    revenue = float((roas_series * data.get("花費金額 (TWD)", pd.Series(dtype=float))).sum())
    conversions = float(data.get("購買次數", pd.Series(dtype=float)).sum())
    impressions = float(data.get("曝光次數", pd.Series(dtype=float)).sum())
    clicks = float(data.get("連結點擊次數", pd.Series(dtype=float)).sum())

    avg_roas = revenue / spend if spend > 0 else 0.0
    cpa = spend / conversions if conversions > 0 else 0.0
    ctr = clicks / impressions * 100 if impressions > 0 else 0.0

    return BaselineMetrics(
        spend=spend,
        revenue=revenue,
        roas=avg_roas,
        conversions=conversions,
        cpa=cpa,
        impressions=impressions,
        clicks=clicks,
        ctr=ctr,
    )


def simulate_roi(
    baseline: BaselineMetrics,
    new_budget: Optional[float] = None,
    budget_change_pct: float = 0.0,
    expected_roas: Optional[float] = None,
    conversion_uplift_pct: float = 0.0,
    fixed_cost: float = 0.0,
    ai_tool_cost: float = 0.0,
) -> Dict[str, Any]:
    baseline_profit = baseline.revenue - baseline.spend

    if new_budget is None:
        new_budget = baseline.spend * (1 + budget_change_pct / 100)
    projected_roas = expected_roas if expected_roas is not None else baseline.roas
    projected_revenue = new_budget * projected_roas
    projected_conversions = baseline.conversions * (1 + conversion_uplift_pct / 100)
    projected_cpa = new_budget / projected_conversions if projected_conversions > 0 else 0.0
    projected_ctr = baseline.ctr * (1 + conversion_uplift_pct / 100)

    gross_profit = projected_revenue - new_budget
    net_profit = gross_profit - fixed_cost - ai_tool_cost
    incremental_revenue = projected_revenue - baseline.revenue
    incremental_spend = new_budget - baseline.spend
    incremental_profit = net_profit - baseline_profit
    roi_ratio = (net_profit / new_budget) if new_budget > 0 else 0.0

    return {
        "baseline": {
            "spend": baseline.spend,
            "revenue": baseline.revenue,
            "roas": baseline.roas,
            "conversions": baseline.conversions,
            "cpa": baseline.cpa,
            "ctr": baseline.ctr,
            "profit": baseline_profit,
        },
        "projected": {
            "spend": new_budget,
            "revenue": projected_revenue,
            "roas": projected_roas,
            "conversions": projected_conversions,
            "cpa": projected_cpa,
            "ctr": projected_ctr,
            "gross_profit": gross_profit,
            "net_profit": net_profit,
            "fixed_cost": fixed_cost,
            "ai_tool_cost": ai_tool_cost,
        },
        "delta": {
            "incremental_spend": incremental_spend,
            "incremental_revenue": incremental_revenue,
            "incremental_profit": incremental_profit,
            "roi": roi_ratio,
        },
    }


def simulation_to_dataframe(result: Dict[str, Any]) -> pd.DataFrame:
    baseline = result["baseline"]
    projected = result["projected"]
    delta = result["delta"]
    data = [
        {
            "指標": "花費",
            "Baseline": baseline["spend"],
            "Projected": projected["spend"],
            "Delta": delta["incremental_spend"],
        },
        {
            "指標": "營收",
            "Baseline": baseline["revenue"],
            "Projected": projected["revenue"],
            "Delta": delta["incremental_revenue"],
        },
        {
            "指標": "ROAS",
            "Baseline": baseline["roas"],
            "Projected": projected["roas"],
            "Delta": projected["roas"] - baseline["roas"],
        },
        {
            "指標": "轉換數",
            "Baseline": baseline["conversions"],
            "Projected": projected["conversions"],
            "Delta": projected["conversions"] - baseline["conversions"],
        },
        {
            "指標": "CPA",
            "Baseline": baseline["cpa"],
            "Projected": projected["cpa"],
            "Delta": projected["cpa"] - baseline["cpa"],
        },
        {
            "指標": "淨利",
            "Baseline": baseline["profit"],
            "Projected": projected["net_profit"],
            "Delta": delta["incremental_profit"],
        },
        {
            "指標": "ROI",
            "Baseline": baseline["profit"] / baseline["spend"] if baseline["spend"] > 0 else 0.0,
            "Projected": delta["roi"],
            "Delta": delta["roi"],
        },
    ]
    return pd.DataFrame(data)


def _format_currency(value: float) -> str:
    return f"NT$ {value:,.0f}"


def _format_float(value: float) -> str:
    return f"{value:.2f}"


def generate_roi_ppt(result: Dict[str, Any], req: ROIReportRequest, output_path: str) -> Path:
    if not PPTX_AVAILABLE:
        raise RuntimeError("尚未安裝 python-pptx，無法匯出 ROI 報表。")

    report = Presentation()
    baseline = result["baseline"]
    projected = result["projected"]
    delta = result["delta"]

    title_slide = report.slide_layouts[0]
    slide = report.slides.add_slide(title_slide)
    slide.shapes.title.text = f"{req.brand_name} ROI 模擬報告"
    subtitle = slide.placeholders[1]
    subtitle.text = f"{req.scenario_name}\n期間：{req.start_date.date()} ～ {req.end_date.date()}"

    summary_slide = report.slide_layouts[1]
    slide = report.slides.add_slide(summary_slide)
    slide.shapes.title.text = "情境摘要"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    summary_points = [
        f"Baseline 花費：{_format_currency(baseline['spend'])}",
        f"Baseline ROAS：{_format_float(baseline['roas'])}",
        f"Projected 花費：{_format_currency(projected['spend'])}",
        f"Projected 淨利：{_format_currency(projected['net_profit'])}",
        f"ROI：{delta['roi'] * 100:.1f}%",
    ]
    for point in summary_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

    table_slide = report.slide_layouts[5]
    slide = report.slides.add_slide(table_slide)
    slide.shapes.title.text = "Baseline vs Projected"
    df = simulation_to_dataframe(result)
    rows, cols = df.shape
    table = slide.shapes.add_table(
        rows + 1,
        cols,
        Inches(0.6),
        Inches(1.6),
        Inches(8.5),
        Inches(3.6),
    ).table

    for idx, column in enumerate(df.columns):
        table.cell(0, idx).text = column
        table.cell(0, idx).text_frame.paragraphs[0].font.bold = True
    for r in range(rows):
        for c in range(cols):
            value = df.iloc[r, c]
            if isinstance(value, (int, float)):
                if df.columns[c] in {"Baseline", "Projected", "Delta"} and df.iloc[r]["指標"] not in {"ROI"}:
                    text = _format_currency(float(value))
                elif df.iloc[r]["指標"] == "ROI":
                    text = f"{float(value) * 100:.1f}%"
                else:
                    text = f"{float(value):.2f}"
            else:
                text = str(value)
            table.cell(r + 1, c).text = text

    assumption_slide = report.slide_layouts[1]
    slide = report.slides.add_slide(assumption_slide)
    slide.shapes.title.text = "假設與備註"
    body = slide.placeholders[1].text_frame
    body.clear()
    assumptions = [
        "Baseline 指標以選定期間歷史資料計算。",
        "Projected ROAS 與轉換提升由使用者提供假設。",
        "淨利 = 模擬營收 − 模擬花費 − 固定成本 − AI 工具成本。",
        "ROI = 淨利 / 模擬花費。",
    ]
    for text in assumptions:
        p = body.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
    if req.notes:
        notes_slide = report.slide_layouts[1]
        slide = report.slides.add_slide(notes_slide)
        slide.shapes.title.text = "其他備註"
        body = slide.placeholders[1].text_frame
        body.clear()
        for note in req.notes:
            p = body.add_paragraph()
            p.text = note
            p.font.size = Pt(16)

    output = Path(output_path).expanduser()
    output.parent.mkdir(parents=True, exist_ok=True)
    report.save(output)
    return output
