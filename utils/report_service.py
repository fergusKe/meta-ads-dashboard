from __future__ import annotations

import os
from dataclasses import dataclass
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict, List, Literal, Optional

import pandas as pd

try:
    from ics import Calendar, Event

    ICS_AVAILABLE = True
except ImportError:
    Calendar = Event = None  # type: ignore
    ICS_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


DATE_FMT = "%Y-%m-%d"


@dataclass
class WeeklyReportRequest:
    brand_name: str
    start_date: datetime
    end_date: datetime
    logo_path: Optional[str] = None
    analyst: str | None = None


def _ensure_output_dir(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def _filter_by_date(df: pd.DataFrame, start: datetime, end: datetime) -> pd.DataFrame:
    if "開始" not in df.columns:
        return df
    mask = (df["開始"] >= pd.Timestamp(start)) & (df["開始"] <= pd.Timestamp(end))
    filtered = df[mask]
    if filtered.empty:
        return df
    return filtered


def run_weekly_qa(df: pd.DataFrame, start: datetime, end: datetime) -> Dict[str, object]:
    """
    執行週報生成前的資料品質檢查，回傳包含狀態、議題與覆蓋率的報告。
    """
    requested_start = start.date()
    requested_end = end.date()
    expected_days = max((requested_end - requested_start).days + 1, 1)

    coverage = {
        "requested_start": requested_start.isoformat(),
        "requested_end": requested_end.isoformat(),
        "data_start": None,
        "data_end": None,
        "days_expected": expected_days,
        "days_with_data": 0,
        "coverage_ratio": 0.0,
    }

    issues: List[Dict[str, str]] = []
    filtered = _filter_by_date(df, start, end).copy()

    if filtered.empty:
        issues.append(
            {
                "level": "error",
                "message": "指定期間沒有投放紀錄，請確認日期區間或資料是否完成更新。",
            }
        )
        return {
            "status": "error",
            "record_count": 0,
            "issues": issues,
            "coverage": coverage,
            "metrics": {},
        }

    if "開始" in filtered.columns:
        start_series = filtered["開始"].dropna()
        if not start_series.empty:
            data_start = start_series.min().date()
            data_end = start_series.max().date()
            coverage["data_start"] = data_start.isoformat()
            coverage["data_end"] = data_end.isoformat()
            unique_days = pd.Series(start_series.dt.date.unique())
            coverage["days_with_data"] = int(unique_days.size)
            coverage["coverage_ratio"] = round(coverage["days_with_data"] / expected_days, 2)
            if coverage["coverage_ratio"] < 0.5:
                issues.append(
                    {
                        "level": "warning",
                        "message": "資料涵蓋天數低於 50%，建議確認是否選擇正確期間或資料是否完整上傳。",
                    }
                )
        else:
            issues.append(
                {
                    "level": "warning",
                    "message": "篩選後的資料缺少「開始」欄位值，無法評估涵蓋天數。",
                }
            )
    else:
        issues.append(
            {
                "level": "error",
                "message": "缺少「開始」欄位，無法依期間篩選資料。",
            }
        )

    required_columns = {
        "花費金額 (TWD)": "廣告花費",
        "購買次數": "購買次數",
        "購買 ROAS（廣告投資報酬率）": "購買 ROAS",
        "曝光次數": "曝光次數",
    }

    for column, label in required_columns.items():
        if column not in filtered.columns:
            issues.append(
                {
                    "level": "error",
                    "message": f"缺少關鍵欄位：{label}（{column}）。",
                }
            )

    metrics = prepare_weekly_metrics(filtered, start, end)
    totals = {
        "total_spend": round(metrics.get("total_spend", 0.0), 2),
        "total_conversions": round(metrics.get("total_conversions", 0.0), 2),
        "weighted_roas": round(metrics.get("weighted_roas", 0.0), 2),
        "ctr": round(metrics.get("ctr", 0.0), 2),
    }

    if totals["total_spend"] <= 0:
        issues.append(
            {
                "level": "warning",
                "message": "本期間總花費為 0，請確認是否尚未投放或資料是否正確匯入。",
            }
        )

    if totals["total_conversions"] <= 0:
        issues.append(
            {
                "level": "warning",
                "message": "本期間無任何轉換紀錄，報告內容可能缺乏成效亮點。",
            }
        )

    status: Literal["ok", "warning", "error"] = "ok"
    if any(issue["level"] == "error" for issue in issues):
        status = "error"
    elif any(issue["level"] == "warning" for issue in issues):
        status = "warning"

    return {
        "status": status,
        "record_count": int(len(filtered)),
        "issues": issues,
        "coverage": coverage,
        "metrics": totals,
    }


def prepare_weekly_metrics(df: pd.DataFrame, start: datetime, end: datetime) -> dict:
    data = _filter_by_date(df, start, end).copy()
    total_spend = float(data.get("花費金額 (TWD)", pd.Series(dtype=float)).sum())
    total_conversions = float(data.get("購買次數", pd.Series(dtype=float)).sum())
    total_impressions = float(data.get("曝光次數", pd.Series(dtype=float)).sum())
    total_clicks = float(data.get("連結點擊次數", pd.Series(dtype=float)).sum())

    weighted_roas = 0.0
    if total_spend > 0 and "購買 ROAS（廣告投資報酬率）" in data.columns:
        weighted_roas = float(
            (data["購買 ROAS（廣告投資報酬率）"] * data["花費金額 (TWD)"]).sum() / total_spend
        )

    ctr = 0.0
    if total_impressions > 0 and total_clicks > 0:
        ctr = total_clicks / total_impressions * 100

    top_campaigns = pd.DataFrame()
    if "行銷活動名稱" in data.columns:
        top_campaigns = (
            data.groupby("行銷活動名稱", as_index=False)
            .agg(
                {
                    "花費金額 (TWD)": "sum",
                    "購買次數": "sum",
                    "購買 ROAS（廣告投資報酬率）": "mean",
                }
            )
            .sort_values("花費金額 (TWD)", ascending=False)
            .head(5)
        )

    top_ads = pd.DataFrame()
    if "廣告名稱" in data.columns:
        top_ads = (
            data.groupby("廣告名稱", as_index=False)
            .agg(
                {
                    "花費金額 (TWD)": "sum",
                    "購買次數": "sum",
                    "購買 ROAS（廣告投資報酬率）": "mean",
                }
            )
            .sort_values("購買次數", ascending=False)
            .head(5)
        )

    return {
        "total_spend": total_spend,
        "total_conversions": total_conversions,
        "weighted_roas": weighted_roas,
        "ctr": ctr,
        "top_campaigns": top_campaigns,
        "top_ads": top_ads,
    }


def _add_title_slide(prs: Presentation, req: WeeklyReportRequest, metrics: dict) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = f"{req.brand_name} AI 投放週報"
    subtitle.text = f"期間：{req.start_date.strftime(DATE_FMT)} ～ {req.end_date.strftime(DATE_FMT)}"

    if req.logo_path and Path(req.logo_path).exists():
        left = Inches(8.0)
        top = Inches(0.2)
        height = Inches(0.8)
        slide.shapes.add_picture(req.logo_path, left, top, height=height)


def _add_kpi_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "關鍵指標摘要"

    body = slide.placeholders[1].text_frame
    body.clear()
    body.word_wrap = True

    spend = metrics["total_spend"]
    roas = metrics["weighted_roas"]
    conversions = metrics["total_conversions"]
    ctr = metrics["ctr"]

    bullet_texts = [
        f"總花費：NT$ {spend:,.0f}",
        f"平均 ROAS：{roas:.2f}",
        f"總轉換次數：{conversions:,.0f}",
        f"平均 CTR：{ctr:.2f}%",
    ]
    for text in bullet_texts:
        p = body.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)


def _add_table_slide(prs: Presentation, title_text: str, df: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text

    if df.empty:
        body = slide.shapes.placeholders[1].text_frame
        body.text = "本週無資料"
        return

    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.8), Inches(9), Inches(3.5)).table

    for idx, col_name in enumerate(df.columns):
        table.cell(0, idx).text = col_name
        table.cell(0, idx).text_frame.paragraphs[0].font.bold = True
        table.cell(0, idx).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for r in range(rows):
        for c in range(cols):
            value = df.iloc[r, c]
            if isinstance(value, (int, float)):
                text = f"{value:,.2f}" if abs(value) < 1000 else f"{value:,.0f}"
            else:
                text = str(value)
            table.cell(r + 1, c).text = text
            table.cell(r + 1, c).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _add_next_steps_slide(prs: Presentation, suggestions: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "下週建議行動"
    body = slide.placeholders[1].text_frame
    body.clear()
    for item in suggestions:
        p = body.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)


def generate_weekly_report(df: pd.DataFrame, req: WeeklyReportRequest, output_path: str, suggestions: Optional[list[str]] = None) -> Path:
    if not PPTX_AVAILABLE:
        raise RuntimeError("尚未安裝 python-pptx，無法產生報表。請先在環境中安裝 python-pptx。")

    metrics = prepare_weekly_metrics(df, req.start_date, req.end_date)
    presentation = Presentation()

    _add_title_slide(presentation, req, metrics)
    _add_kpi_slide(presentation, metrics)

    if not metrics["top_campaigns"].empty:
        _add_table_slide(
            presentation,
            "花費最高活動 Top 5",
            metrics["top_campaigns"].rename(
                columns={
                    "花費金額 (TWD)": "花費",
                    "購買次數": "轉換",
                    "購買 ROAS（廣告投資報酬率）": "平均 ROAS",
                }
            ),
        )

    if not metrics["top_ads"].empty:
        _add_table_slide(
            presentation,
            "轉換最佳素材 Top 5",
            metrics["top_ads"].rename(
                columns={
                    "花費金額 (TWD)": "花費",
                    "購買次數": "轉換",
                    "購買 ROAS（廣告投資報酬率）": "平均 ROAS",
                }
            ),
        )

    _add_next_steps_slide(
        presentation,
        suggestions
        or [
            "檢視 ROAS 低於 2.0 的活動，評估是否停掉或調整受眾。",
            "延伸轉換最佳素材，進行變體測試。",
            "盤點下一檔檔期素材需求，提早準備文案與圖像。",
        ],
    )

    output = Path(output_path).expanduser()
    _ensure_output_dir(output)
    presentation.save(output)
    return output


def generate_monthly_summary(df: pd.DataFrame, req: WeeklyReportRequest) -> dict:
    def _slice_by_period(source: pd.DataFrame, start_dt: datetime, end_dt: datetime) -> pd.DataFrame:
        if "開始" not in source.columns:
            return source.copy()
        mask = (source["開始"] >= pd.Timestamp(start_dt)) & (source["開始"] <= pd.Timestamp(end_dt))
        return source.loc[mask].copy()

    def _aggregate_metrics(data: pd.DataFrame) -> Dict[str, float]:
        if data.empty:
            return {
                "total_spend": 0.0,
                "total_conversions": 0.0,
                "revenue": 0.0,
                "weighted_roas": 0.0,
                "ctr": 0.0,
                "cpa": 0.0,
            }
        spend = float(data.get("花費金額 (TWD)", pd.Series(dtype=float)).sum())
        roas_series = data.get("購買 ROAS（廣告投資報酬率）", pd.Series(dtype=float))
        revenue = float((roas_series * data.get("花費金額 (TWD)", pd.Series(dtype=float))).sum()) if not data.empty else 0.0
        conversions = float(data.get("購買次數", pd.Series(dtype=float)).sum())
        impressions = float(data.get("曝光次數", pd.Series(dtype=float)).sum())
        clicks = float(data.get("連結點擊次數", pd.Series(dtype=float)).sum())

        weighted_roas = revenue / spend if spend > 0 else 0.0
        ctr = clicks / impressions * 100 if impressions > 0 else 0.0
        cpa = spend / conversions if conversions > 0 else 0.0

        return {
            "total_spend": spend,
            "total_conversions": conversions,
            "revenue": revenue,
            "weighted_roas": weighted_roas,
            "ctr": ctr,
            "cpa": cpa,
        }

    def _top_campaigns(data: pd.DataFrame, limit: int = 5) -> pd.DataFrame:
        if data.empty or "行銷活動名稱" not in data.columns:
            return pd.DataFrame()
        return (
            data.groupby("行銷活動名稱", as_index=False)
            .agg(
                {
                    "花費金額 (TWD)": "sum",
                    "購買次數": "sum",
                    "購買 ROAS（廣告投資報酬率）": "mean",
                }
            )
            .sort_values("花費金額 (TWD)", ascending=False)
            .head(limit)
        )

    def _lagging_campaigns(data: pd.DataFrame, limit: int = 5, roas_threshold: float = 1.0) -> pd.DataFrame:
        if data.empty or "行銷活動名稱" not in data.columns or "購買 ROAS（廣告投資報酬率）" not in data.columns:
            return pd.DataFrame()
        filtered = data[data["購買 ROAS（廣告投資報酬率）"] < roas_threshold]
        if filtered.empty:
            return pd.DataFrame()
        return (
            filtered.groupby("行銷活動名稱", as_index=False)
            .agg(
                {
                    "花費金額 (TWD)": "sum",
                    "購買次數": "sum",
                    "購買 ROAS（廣告投資報酬率）": "mean",
                }
            )
            .sort_values("花費金額 (TWD)", ascending=False)
            .head(limit)
        )

    def _channel_breakdown(data: pd.DataFrame, limit: int = 10) -> pd.DataFrame:
        if data.empty or "行銷活動名稱" not in data.columns:
            return pd.DataFrame()
        agg = (
            data.groupby("行銷活動名稱", as_index=False)[["花費金額 (TWD)", "購買次數"]]
            .sum()
            .sort_values("花費金額 (TWD)", ascending=False)
            .head(limit)
        )
        return agg

    def _pct_change(current: float, previous: float) -> float | None:
        if previous == 0:
            return None
        return (current - previous) / previous * 100

    current_data = _slice_by_period(df, req.start_date, req.end_date)
    period_length = max((req.end_date.date() - req.start_date.date()).days + 1, 1)
    prev_end_date = req.start_date.date() - timedelta(days=1)
    prev_start_date = prev_end_date - timedelta(days=period_length - 1)
    prev_start_dt = datetime.combine(prev_start_date, datetime.min.time())
    prev_end_dt = datetime.combine(prev_end_date, datetime.max.time())
    previous_data = _slice_by_period(df, prev_start_dt, prev_end_dt) if prev_end_date >= prev_start_date else pd.DataFrame()

    current_metrics = _aggregate_metrics(current_data)
    previous_metrics = _aggregate_metrics(previous_data)

    top_campaigns_df = _top_campaigns(current_data)
    lagging_campaigns_df = _lagging_campaigns(current_data)
    channel_breakdown_df = _channel_breakdown(current_data)

    delta = {
        "spend": current_metrics["total_spend"] - previous_metrics["total_spend"],
        "spend_pct": _pct_change(current_metrics["total_spend"], previous_metrics["total_spend"]),
        "conversions": current_metrics["total_conversions"] - previous_metrics["total_conversions"],
        "conversions_pct": _pct_change(current_metrics["total_conversions"], previous_metrics["total_conversions"]),
        "roas": current_metrics["weighted_roas"] - previous_metrics["weighted_roas"],
        "roas_pct": _pct_change(current_metrics["weighted_roas"], previous_metrics["weighted_roas"]),
        "ctr": current_metrics["ctr"] - previous_metrics["ctr"],
        "ctr_pct": _pct_change(current_metrics["ctr"], previous_metrics["ctr"]),
    }

    highlights: List[str] = []
    warnings: List[str] = []

    if not top_campaigns_df.empty:
        top_row = top_campaigns_df.iloc[0]
        highlights.append(
            f"活動「{top_row['行銷活動名稱']}」貢獻 NT$ {top_row['花費金額 (TWD)']:,.0f} 花費與 {top_row['購買次數']:.0f} 次轉換，平均 ROAS {top_row['購買 ROAS（廣告投資報酬率）']:.2f}。"
        )

    if delta["conversions_pct"] is not None and delta["conversions_pct"] > 0:
        highlights.append(f"整體轉換次數較前期成長 {delta['conversions_pct']:.1f}%。")

    if delta["roas"] > 0:
        highlights.append(f"平均 ROAS 提升至 {current_metrics['weighted_roas']:.2f}（較前期提高 {delta['roas']:.2f}）。")

    if lagging_campaigns_df.empty and current_metrics["total_spend"] > 0:
        highlights.append("所有活動 ROAS 均高於 1.0，持續維持高效投放。")

    if not lagging_campaigns_df.empty:
        bottom_row = lagging_campaigns_df.iloc[0]
        warnings.append(
            f"活動「{bottom_row['行銷活動名稱']}」ROAS 僅 {bottom_row['購買 ROAS（廣告投資報酬率）']:.2f}，本月仍投入 NT$ {bottom_row['花費金額 (TWD)']:,.0f}，建議檢視素材與受眾設定。"
        )

    if delta["spend_pct"] is not None and delta["spend_pct"] > 20:
        warnings.append(f"整體花費較前期增加 {delta['spend_pct']:.1f}%，請留意預算控管與成效。")

    summary = {
        "brand_name": req.brand_name,
        "period": {
            "current_start": req.start_date.date().isoformat(),
            "current_end": req.end_date.date().isoformat(),
            "previous_start": prev_start_dt.date().isoformat(),
            "previous_end": prev_end_dt.date().isoformat(),
            "days": period_length,
        },
        "summary_metrics": current_metrics,
        "trend": {
            "previous_metrics": previous_metrics,
            "delta": delta,
        },
        "highlights": highlights,
        "warnings": warnings,
        "top_campaigns": top_campaigns_df.to_dict(orient="records"),
        "lagging_campaigns": lagging_campaigns_df.to_dict(orient="records"),
        "channel_breakdown": channel_breakdown_df.to_dict(orient="records"),
    }

    return summary


def build_monthly_email_html(summary: Dict[str, object]) -> str:
    brand = summary.get("brand_name", "品牌")
    period = summary.get("period", {})
    metrics: Dict[str, float] = summary.get("summary_metrics", {}) or {}
    trend: Dict[str, Dict[str, float]] = summary.get("trend", {}) or {}
    delta = trend.get("delta", {}) or {}
    previous_metrics = trend.get("previous_metrics", {}) or {}
    highlights: List[str] = summary.get("highlights", []) or []
    warnings: List[str] = summary.get("warnings", []) or []
    custom_suggestions: List[str] = summary.get("custom_suggestions", []) or []

    def _format_ratio(value: float | None) -> str:
        if value is None:
            return "—"
        sign = "+" if value >= 0 else ""
        return f"{sign}{value:.1f}%"

    def _format_currency(value: float | None) -> str:
        return f"NT$ {value:,.0f}" if value else "NT$ 0"

    def _format_number(value: float | None) -> str:
        return f"{value:,.0f}" if value else "0"

    def _format_float(value: float | None) -> str:
        return f"{value:.2f}" if value else "0.00"

    def _render_table(records: List[Dict[str, object]], rename: Dict[str, str]) -> str:
        if not records:
            return "<p style='color:#718096;'>本月無資料。</p>"
        df = pd.DataFrame(records).rename(columns=rename)
        for col in df.columns:
            if "NT$" in col or "花費" in col:
                df[col] = df[col].map(lambda x: f"{float(x):,.0f}" if pd.notna(x) else "0")
            elif df[col].dtype.kind in {"f", "c"}:
                df[col] = df[col].map(lambda x: f"{float(x):.2f}" if pd.notna(x) else "0.00")
        return df.to_html(index=False, escape=False, classes="table", justify="center")

    channel_html = _render_table(
        summary.get("channel_breakdown", []) or [],
        rename={"行銷活動名稱": "活動", "花費金額 (TWD)": "花費 (NT$)", "購買次數": "轉換"},
    )
    top_campaigns_html = _render_table(
        summary.get("top_campaigns", []) or [],
        rename={
            "行銷活動名稱": "活動",
            "花費金額 (TWD)": "花費 (NT$)",
            "購買次數": "轉換",
            "購買 ROAS（廣告投資報酬率）": "平均 ROAS",
        },
    )
    lagging_campaigns_html = _render_table(
        summary.get("lagging_campaigns", []) or [],
        rename={
            "行銷活動名稱": "活動",
            "花費金額 (TWD)": "花費 (NT$)",
            "購買次數": "轉換",
            "購買 ROAS（廣告投資報酬率）": "平均 ROAS",
        },
    )

    custom_note_html = ""
    if custom_suggestions:
        custom_note_html = (
            "<div class='card'>"
            "<h3>客製寄送備註</h3>"
            + ("<ul>" + "".join(f"<li>{item}</li>" for item in custom_suggestions) + "</ul>")
            + "</div>"
        )

    html = f"""
<!DOCTYPE html>
<html lang="zh-Hant">
<head>
  <meta charset="UTF-8" />
  <style>
    body {{
      font-family: "PingFang TC", "Microsoft JhengHei", Arial, sans-serif;
      color: #1f2933;
      margin: 0;
      padding: 24px;
      background-color: #f7fafc;
    }}
    h2 {{
      margin-top: 0;
      color: #1a365d;
    }}
    .card {{
      background: #ffffff;
      padding: 20px;
      border-radius: 12px;
      box-shadow: 0 8px 20px rgba(15, 23, 42, 0.08);
      margin-bottom: 18px;
    }}
    .metrics-grid {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(160px, 1fr));
      gap: 12px;
    }}
    .metric {{
      background: #f1f5f9;
      border-radius: 8px;
      padding: 12px;
    }}
    .metric span {{
      display: block;
    }}
    .metric .label {{
      color: #475569;
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: 0.2em;
    }}
    .metric .value {{
      font-size: 20px;
      font-weight: 600;
      color: #1e293b;
    }}
    .metric .delta {{
      font-size: 12px;
      color: #047857;
    }}
    .metric .delta.negative {{
      color: #b91c1c;
    }}
    ul {{
      padding-left: 20px;
      margin: 8px 0;
    }}
    .table {{
      width: 100%;
      border-collapse: collapse;
      margin-top: 12px;
    }}
    .table th, .table td {{
      border: 1px solid #e2e8f0;
      padding: 8px;
      font-size: 14px;
      text-align: center;
    }}
    .table th {{
      background-color: #edf2f7;
      color: #1a202c;
    }}
    .note {{
      font-size: 12px;
      color: #64748b;
    }}
  </style>
</head>
<body>
  <h2>{brand}｜{period.get("current_start", "-")} ～ {period.get("current_end", "-")} 月報摘要</h2>

  <div class="card">
    <div class="metrics-grid">
      <div class="metric">
        <span class="label">總花費</span>
        <span class="value">{_format_currency(metrics.get("total_spend"))}</span>
        <span class="delta{' negative' if (delta.get('spend') or 0) < 0 else ''}">vs 前期 {_format_ratio(delta.get('spend_pct'))}</span>
      </div>
      <div class="metric">
        <span class="label">轉換次數</span>
        <span class="value">{_format_number(metrics.get("total_conversions"))}</span>
        <span class="delta{' negative' if (delta.get('conversions') or 0) < 0 else ''}">vs 前期 {_format_ratio(delta.get('conversions_pct'))}</span>
      </div>
      <div class="metric">
        <span class="label">平均 ROAS</span>
        <span class="value">{_format_float(metrics.get("weighted_roas"))}</span>
        <span class="delta{' negative' if delta.get('roas') and delta.get('roas') < 0 else ''}">vs 前期 {_format_ratio(delta.get('roas_pct'))}</span>
      </div>
      <div class="metric">
        <span class="label">平均 CTR</span>
        <span class="value">{_format_float(metrics.get("ctr"))}%</span>
        <span class="delta{' negative' if delta.get('ctr') and delta.get('ctr') < 0 else ''}">vs 前期 {_format_ratio(delta.get('ctr_pct'))}</span>
      </div>
    </div>
    <p class="note">前期期間：{period.get("previous_start", "-")} ～ {period.get("previous_end", "-")}</p>
  </div>

  <div class="card">
    <h3>亮點觀察</h3>
    {"<ul>" + "".join(f"<li>{item}</li>" for item in highlights) + "</ul>" if highlights else "<p style='color:#718096;'>本期尚無特別亮點。</p>"}
  </div>

  <div class="card">
    <h3>優化建議</h3>
    {"<ul>" + "".join(f"<li>{item}</li>" for item in warnings) + "</ul>" if warnings else "<p style='color:#718096;'>本期未發現急迫的優化項目。</p>"}
  </div>

  <div class="card">
    <h3>通路 / 活動花費分布</h3>
    {channel_html}
  </div>

  <div class="card">
    <h3>表現最佳活動</h3>
    {top_campaigns_html}
  </div>

  <div class="card">
    <h3>待優化活動</h3>
    {lagging_campaigns_html}
  </div>

  {custom_note_html}

  <p class="note">此報表由 AI 投放助手自動生成，請於寄送前再次複核關鍵數據。</p>
</body>
</html>
"""
    return html


def create_calendar_event(subject: str, start: datetime, end: datetime, description: str = "") -> str:
    if ICS_AVAILABLE:
        calendar = Calendar()
        event = Event()
        event.name = subject
        event.begin = start
        event.end = end
        event.description = description
        calendar.events.add(event)
        return str(calendar)

    def _format_datetime(dt: datetime) -> str:
        return dt.strftime("%Y%m%dT%H%M%S")

    lines = [
        "BEGIN:VCALENDAR",
        "VERSION:2.0",
        "PRODID:-//AI Report Scheduler//EN",
        "BEGIN:VEVENT",
        f"SUMMARY:{subject}",
        f"DTSTART:{_format_datetime(start)}",
        f"DTEND:{_format_datetime(end)}",
    ]
    if description:
        sanitized = description.replace("\n", "\\n")
        lines.append(f"DESCRIPTION:{sanitized}")
    lines.extend(["END:VEVENT", "END:VCALENDAR"])
    return "\n".join(lines)
